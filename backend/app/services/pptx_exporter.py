from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from app.core.config import settings
from app.schemas.presentation import FinalPresentation


SLIDE_W = 13.333
SLIDE_H = 7.5


THEMES = {
    "light": {
        "bg": "F8FAFC",
        "text": "0F172A",
        "muted_text": "475569",
        "card": "FFFFFF",
        "card_alt": "F1F5F9",
        "accent": "6366F1",
        "accent_soft": "E0E7FF",
        "line": "CBD5E1",
    },
    "accent": {
        "bg": "ECFDF5",
        "text": "064E3B",
        "muted_text": "166534",
        "card": "FFFFFF",
        "card_alt": "D1FAE5",
        "accent": "10B981",
        "accent_soft": "A7F3D0",
        "line": "6EE7B7",
    },
    "muted": {
        "bg": "F5F3FF",
        "text": "3B0764",
        "muted_text": "6B21A8",
        "card": "FFFFFF",
        "card_alt": "EDE9FE",
        "accent": "8B5CF6",
        "accent_soft": "DDD6FE",
        "line": "C4B5FD",
    },
    "white": {
        "bg": "FFFFFF",
        "text": "111827",
        "muted_text": "4B5563",
        "card": "F9FAFB",
        "card_alt": "F3F4F6",
        "accent": "2563EB",
        "accent_soft": "DBEAFE",
        "line": "D1D5DB",
    },
    "dark": {
        "bg": "EEF2FF",
        "text": "1E1B4B",
        "muted_text": "4338CA",
        "card": "FFFFFF",
        "card_alt": "E0E7FF",
        "accent": "6366F1",
        "accent_soft": "C7D2FE",
        "line": "A5B4FC",
    },
}


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.replace("#", "").strip())


def _theme(style: str | None):
    return THEMES.get(style or "light", THEMES["light"])


def _set_fill(shape, color: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)


def _set_line(shape, color: str, width_pt: float = 1.0):
    shape.line.color.rgb = _rgb(color)
    shape.line.width = Pt(width_pt)


def _apply_bg(slide, theme_key: str | None):
    t = _theme(theme_key)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(t["bg"])


def _add_top_bar(slide, theme_key: str | None):
    t = _theme(theme_key)
    top_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(SLIDE_W),
        Inches(0.14),
    )
    _set_fill(top_bar, t["accent"])
    top_bar.line.fill.background()


def _add_card(slide, left, top, width, height, theme_key: str | None, alt: bool = False):
    t = _theme(theme_key)
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    _set_fill(shape, t["card_alt"] if alt else t["card"])
    _set_line(shape, t["line"], 1)
    return shape


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    theme_key: str | None,
    font_size=16,
    bold=False,
    color=None,
    align=PP_ALIGN.LEFT,
):
    t = _theme(theme_key)
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color or t["text"])
    p.alignment = align
    return tx


def _fit_title(slide, text: str, theme_key: str | None, max_size=28, min_size=18):
    t = _theme(theme_key)
    tx = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(11.8), Inches(1.0))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.color.rgb = _rgb(t["text"])

    if len(text) > 80:
        p.font.size = Pt(min_size)
    elif len(text) > 50:
        p.font.size = Pt(22)
    else:
        p.font.size = Pt(max_size)

    return tx


def _column_labels(item):
    bullets = item.bullets or []
    left_label = "Ключевые идеи"
    right_label = "Практические аспекты"

    if len(bullets) >= 2:
        left_label = bullets[0][:32].strip("•-–—:;,. ")
        right_label = bullets[1][:32].strip("•-–—:;,. ")

        if len(left_label) < 4:
            left_label = "Ключевые идеи"
        if len(right_label) < 4:
            right_label = "Практические аспекты"

    return left_label, right_label


def _insert_image(slide, image_path: str, left, top, width, height):
    if image_path and Path(image_path).exists():
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))


def _render_title_and_bullets(slide, item):
    theme_key = getattr(item, "background_style", "light")
    t = _theme(theme_key)

    _add_card(slide, 0.7, 1.65, 11.9, 4.95, theme_key, alt=False)

    bullets = item.bullets[:6] if item.bullets else ["Ключевая мысль по теме слайда"]

    y = 2.0
    for bullet in bullets:
        accent = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(1.0),
            Inches(y + 0.09),
            Inches(0.16),
            Inches(0.16),
        )
        _set_fill(accent, t["accent"])
        accent.line.fill.background()

        _add_textbox(slide, 1.35, y - 0.05, 10.4, 0.55, bullet, theme_key, font_size=18)
        y += 0.7


def _render_two_columns(slide, item):
    theme_key = getattr(item, "background_style", "light")
    bullets = item.bullets[:8] if item.bullets else ["Пункт 1", "Пункт 2", "Пункт 3", "Пункт 4"]
    mid = max(2, len(bullets) // 2)

    left_col = bullets[:mid]
    right_col = bullets[mid:]
    left_label, right_label = _column_labels(item)

    _add_card(slide, 0.75, 1.75, 5.75, 4.9, theme_key, alt=False)
    _add_card(slide, 6.8, 1.75, 5.75, 4.9, theme_key, alt=True)

    _add_textbox(slide, 1.05, 2.0, 4.95, 0.4, left_label, theme_key, font_size=15, bold=True)
    _add_textbox(slide, 7.1, 2.0, 4.95, 0.4, right_label, theme_key, font_size=15, bold=True)

    y = 2.45
    for bullet in left_col:
        _add_textbox(slide, 1.05, y, 4.95, 0.5, f"• {bullet}", theme_key, font_size=17)
        y += 0.62

    y = 2.45
    for bullet in right_col:
        _add_textbox(slide, 7.1, y, 4.95, 0.5, f"• {bullet}", theme_key, font_size=17)
        y += 0.62


def _render_comparison(slide, item):
    theme_key = getattr(item, "background_style", "light")
    t = _theme(theme_key)
    bullets = item.bullets[:8] if item.bullets else ["Фактор 1", "Фактор 2", "Фактор 3", "Фактор 4"]
    left = bullets[::2]
    right = bullets[1::2] or ["Дополнительный аспект"]

    _add_card(slide, 0.75, 1.75, 5.75, 4.95, theme_key, alt=False)
    _add_card(slide, 6.8, 1.75, 5.75, 4.95, theme_key, alt=True)

    _add_textbox(slide, 1.05, 2.0, 4.95, 0.35, "Подход A", theme_key, font_size=16, bold=True, color=t["accent"])
    _add_textbox(slide, 7.1, 2.0, 4.95, 0.35, "Подход B", theme_key, font_size=16, bold=True, color=t["accent"])

    y = 2.5
    for bullet in left:
        _add_textbox(slide, 1.05, y, 4.95, 0.52, bullet, theme_key, font_size=17)
        y += 0.62

    y = 2.5
    for bullet in right:
        _add_textbox(slide, 7.1, y, 4.95, 0.52, bullet, theme_key, font_size=17)
        y += 0.62


def _render_timeline(slide, item):
    theme_key = getattr(item, "background_style", "light")
    t = _theme(theme_key)
    bullets = item.bullets[:4] if item.bullets else ["Этап 1", "Этап 2", "Этап 3"]
    n = len(bullets)

    _add_card(slide, 0.7, 1.85, 11.95, 4.5, theme_key, alt=False)

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1.45),
        Inches(4.1),
        Inches(10.1),
        Inches(0.08),
    )
    _set_fill(line, t["accent"])
    line.line.fill.background()

    if n == 1:
        positions = [1.45 + 10.1 / 2]
    else:
        step = 10.1 / (n - 1)
        positions = [1.45 + i * step for i in range(n)]

    for i, (x, bullet) in enumerate(zip(positions, bullets), start=1):
        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x - 0.25),
            Inches(3.82),
            Inches(0.5),
            Inches(0.5),
        )
        _set_fill(circle, t["accent"])
        circle.line.fill.background()
        circle.text_frame.text = str(i)
        circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        circle.text_frame.paragraphs[0].font.color.rgb = _rgb("FFFFFF")
        circle.text_frame.paragraphs[0].font.bold = True
        circle.text_frame.paragraphs[0].font.size = Pt(13)

        box_w = 2.25
        left = max(0.85, min(x - box_w / 2, SLIDE_W - box_w - 0.85))

        _add_textbox(
            slide,
            left,
            4.55,
            box_w,
            0.95,
            bullet,
            theme_key,
            font_size=14,
            align=PP_ALIGN.CENTER,
        )


def _render_metrics(slide, item):
    theme_key = getattr(item, "background_style", "light")
    bullets = item.bullets[:6] if item.bullets else [
        "Метрика 1",
        "Метрика 2",
        "Метрика 3",
        "Метрика 4",
    ]

    positions = [
        (0.85, 1.95),
        (6.85, 1.95),
        (0.85, 3.65),
        (6.85, 3.65),
        (0.85, 5.35),
        (6.85, 5.35),
    ]

    card_h = 1.15

    for idx, ((left, top), bullet) in enumerate(zip(positions, bullets), start=1):
        _add_card(slide, left, top, 5.05, card_h, theme_key, alt=(idx % 2 == 0))
        _add_textbox(
            slide,
            left + 0.22,
            top + 0.18,
            4.6,
            0.72,
            bullet,
            theme_key,
            font_size=16,
            bold=True,
            align=PP_ALIGN.CENTER,
        )


def _render_section(slide, item):
    theme_key = getattr(item, "background_style", "light")
    t = _theme(theme_key)

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1.0),
        Inches(2.35),
        Inches(11.2),
        Inches(2.0),
    )
    _set_fill(band, t["accent_soft"])
    band.line.fill.background()

    subtitle = item.bullets[0] if item.bullets else item.visual_idea or "Раздел презентации"
    _add_textbox(slide, 1.55, 2.78, 10.0, 0.8, subtitle, theme_key, font_size=23, bold=True, align=PP_ALIGN.CENTER)


def _render_image_focus(slide, item):
    theme_key = getattr(item, "background_style", "light")

    has_image = bool(getattr(item, "image_path", None))

    if has_image:
        _insert_image(slide, item.image_path, 0.9, 1.8, 7.2, 4.8)
    else:
        _add_card(slide, 0.9, 1.8, 7.2, 4.8, theme_key, alt=True)
        _add_textbox(
            slide,
            1.35,
            3.7,
            6.2,
            0.6,
            "Иллюстрация недоступна",
            theme_key,
            font_size=18,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    _add_card(slide, 8.45, 1.9, 3.75, 4.55, theme_key, alt=True)
    bullets = item.bullets[:5] if item.bullets else ["Ключевая идея"]

    y = 2.25
    for bullet in bullets:
        _add_textbox(slide, 8.75, y, 3.0, 0.55, bullet, theme_key, font_size=16)
        y += 0.68


def _render_visual_split(slide, item):
    theme_key = getattr(item, "background_style", "light")

    _add_card(slide, 0.75, 1.8, 5.45, 4.75, theme_key, alt=False)

    if getattr(item, "image_path", None):
        _insert_image(slide, item.image_path, 6.55, 1.9, 5.5, 4.55)
    else:
        _add_card(slide, 6.55, 1.9, 5.5, 4.55, theme_key, alt=True)
        _add_textbox(
            slide,
            7.15,
            3.8,
            4.3,
            0.6,
            "Иллюстрация недоступна",
            theme_key,
            font_size=18,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    bullets = item.bullets[:6] if item.bullets else ["Ключевая мысль"]
    y = 2.15
    for bullet in bullets:
        _add_textbox(slide, 1.05, y, 4.7, 0.5, f"• {bullet}", theme_key, font_size=17)
        y += 0.62


def _render_minimal_statement(slide, item):
    theme_key = getattr(item, "background_style", "light")
    t = _theme(theme_key)

    statement = item.bullets[0] if item.bullets else item.visual_idea or "Ключевая мысль"
    _add_textbox(slide, 1.25, 2.1, 10.6, 1.6, statement, theme_key, font_size=28, bold=True, align=PP_ALIGN.CENTER)

    if len(item.bullets) > 1:
        _add_textbox(
            slide,
            2.0,
            4.55,
            9.2,
            0.7,
            " • ".join(item.bullets[1:4]),
            theme_key,
            font_size=16,
            color=t["muted_text"],
            align=PP_ALIGN.CENTER,
        )


def _render_closing(slide, item):
    theme_key = getattr(item, "background_style", "light")
    t = _theme(theme_key)

    _add_card(slide, 1.1, 2.0, 11.1, 3.55, theme_key, alt=True)
    text = item.bullets[0] if item.bullets else "Спасибо за внимание"

    _add_textbox(slide, 2.0, 2.65, 9.2, 0.85, text, theme_key, font_size=25, bold=True, align=PP_ALIGN.CENTER)

    if len(item.bullets) > 1:
        _add_textbox(
            slide,
            2.0,
            4.0,
            9.2,
            0.8,
            " • ".join(item.bullets[1:5]),
            theme_key,
            font_size=15,
            color=t["muted_text"],
            align=PP_ALIGN.CENTER,
        )


def export_presentation(data: FinalPresentation) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    cover = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_bg(cover, "accent")
    _add_top_bar(cover, "accent")

    _add_textbox(cover, 0.95, 1.55, 10.8, 1.3, data.title, "accent", font_size=30, bold=True)
    _add_textbox(
        cover,
        0.98,
        3.0,
        8.2,
        0.8,
        data.subtitle or f"Аудитория: {data.audience}",
        "accent",
        font_size=17,
        color=_theme("accent")["muted_text"],
    )

    for item in data.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        theme_key = getattr(item, "background_style", "light")
        _apply_bg(slide, theme_key)
        _add_top_bar(slide, theme_key)
        _fit_title(slide, item.title, theme_key, max_size=27, min_size=18)

        layout = item.layout or "title_and_bullets"

        if layout == "two_columns":
            _render_two_columns(slide, item)
        elif layout == "comparison":
            _render_comparison(slide, item)
        elif layout == "timeline":
            _render_timeline(slide, item)
        elif layout == "metrics":
            _render_metrics(slide, item)
        elif layout == "section":
            _render_section(slide, item)
        elif layout == "image_focus":
            _render_image_focus(slide, item)
        elif layout == "visual_split":
            _render_visual_split(slide, item)
        elif layout == "minimal_statement":
            _render_minimal_statement(slide, item)
        elif layout == "closing":
            _render_closing(slide, item)
        else:
            _render_title_and_bullets(slide, item)

        if item.speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = item.speaker_notes

    out_dir = Path(settings.generated_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    path = out_dir / filename
    prs.save(str(path))
    return path